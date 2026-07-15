# Builds docs/HealthGrid-Pitch-Final-v2.pptx from the team's approved deck:
# updates stale numbers and inserts two new-feature slides (notification
# center, district PDF report) cloned from the existing "NEW FEATURE" slide
# so the visual language stays identical.
import copy
import io
import sys
from pptx import Presentation
from pptx.util import Inches

SRC = r"C:\Users\nisha\Downloads\Final_HealthGrid-AI (3).pptx"
OUT = sys.argv[1] if len(sys.argv) > 1 else r"docs/HealthGrid-Pitch-Final-v2.pptx"

pres = Presentation(SRC)


def shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise KeyError(name)


def set_runs(shape, para_runs):
    """para_runs: list of lists of strings, matched to existing paragraphs/runs."""
    tf = shape.text_frame
    for pi, runs in enumerate(para_runs):
        for ri, text in enumerate(runs):
            tf.paragraphs[pi].runs[ri].text = text


# ---- Update stale copy on existing slides ----
s5 = pres.slides[4]
for sh in s5.shapes:
    if sh.has_text_frame:
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if "74 unit tests" in run.text:
                    run.text = run.text.replace("74 unit tests", "88 unit tests")

s9 = pres.slides[8]
for sh in s9.shapes:
    if not sh.has_text_frame:
        continue
    for para in sh.text_frame.paragraphs:
        for run in para.runs:
            if "74 unit tests" in run.text:
                run.text = run.text.replace("74 unit tests", "88 unit tests")
            if run.text.startswith("Admin SDK transactions"):
                run.text = "Admin SDK transactions, security rules · deployed on Google Cloud Run"

A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def set_bullets(shape, texts):
    """Replace a bullet list's paragraphs with `texts`, keeping formatting."""
    tf = shape.text_frame
    template = copy.deepcopy(tf.paragraphs[0]._p)
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    first.runs[0].text = texts[0]
    for extra in first.runs[1:]:
        extra.text = ""
    anchor = first._p
    for text in texts[1:]:
        p = copy.deepcopy(template)
        runs = p.findall(f".//{A_NS}t")
        runs[0].text = text
        for r in runs[1:]:
            r.text = ""
        anchor.addnext(p)
        anchor = p


def set_bullets_labeled(shape, pairs):
    """Replace a bold-label bullet list (runs: [label, rest]) with `pairs`."""
    tf = shape.text_frame
    template = copy.deepcopy(tf.paragraphs[0]._p)
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)

    def fill(p_el, label, rest):
        runs = p_el.findall(f".//{A_NS}t")
        runs[0].text = label
        if len(runs) > 1:
            runs[1].text = rest
            for r in runs[2:]:
                r.text = ""
        else:
            runs[0].text = label + rest

    fill(tf.paragraphs[0]._p, *pairs[0])
    anchor = tf.paragraphs[0]._p
    for label, rest in pairs[1:]:
        p = copy.deepcopy(template)
        fill(p, label, rest)
        anchor.addnext(p)
        anchor = p


def set_first_text(slide, current, new):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip().startswith(current):
            sh.text_frame.paragraphs[0].runs[0].text = new
            return sh
    raise KeyError(current)


# ---- Slide 10 "Pilot-ready" becomes the go-to-market slide ----
# The insight judges reward: hospitals never buy software; districts adopt it
# through NHM's existing budget process. Prove -> Fund -> Scale.
s10 = pres.slides[9]
set_first_text(s10, "Deployable now", "Prove — one district")
set_first_text(s10, "Pilot model", "Fund — through NHM")
set_first_text(s10, "Scale path", "Scale — on existing rails")
set_bullets(
    next(sh for sh in s10.shapes if sh.has_text_frame and sh.text_frame.text.startswith("Works as a layer")),
    [
        "Pilot 10–20 PHCs/CHCs — this hackathon's prize is the pilot",
        "No hardware, no new data entry — runs on records districts already keep",
        "Two KPIs: stock-out days prevented · alert-to-action time",
    ],
)
set_bullets(
    next(sh for sh in s10.shapes if sh.has_text_frame and sh.text_frame.text.startswith("Start with 10")),
    [
        "District evidence → state NHM Program Implementation Plan budget line",
        "The District Health Society review already consumes our PDF report",
        "Runs in low thousands ₹/month per district — inside existing IT budgets",
    ],
)
set_bullets(
    next(sh for sh in s10.shapes if sh.has_text_frame and sh.text_frame.text.startswith("District → state")),
    [
        "HMIS/DVDMS integration — our schema mirrors their entities",
        "District → 36 districts of Maharashtra → 800+ nationally",
        "Scaling is configuration, not code: medicines, languages, geography",
    ],
)
set_first_text(
    s10,
    "HealthGrid turns fragmented facility data",
    "We sell evidence, not software — the pilot's own numbers make the case to the state.",
)

# ---- Clone the stress-mode slide as a template for two new slides ----
TEMPLATE_IDX = 6
src = pres.slides[TEMPLATE_IDX]
info_icon_blob = None
for sh in src.shapes:
    if sh.shape_type == 13 and sh.name == "Image 1":
        info_icon_blob = sh.image.blob


def clone_template():
    new = pres.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        if sh.shape_type == 13:  # pictures re-added fresh (avoids rel surgery)
            continue
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    if info_icon_blob:
        new.shapes.add_picture(io.BytesIO(info_icon_blob), Inches(7.06), Inches(5.53), Inches(0.22), Inches(0.17))
    # A touch more breathing room for the bullet list than the template allows.
    t5 = shape_by_name(new, "Text 5")
    t5.top = Inches(3.92)
    t5.height = Inches(1.28)
    return new


# Slide A: Operational Notification Center
a = clone_template()
set_runs(shape_by_name(a, "Text 2"), [["Operational Notification Center"]])
set_runs(shape_by_name(a, "Text 3"), [["From decision to acknowledgement"]])
set_runs(shape_by_name(a, "Text 4"), [[
    "The command centre turns a facility's live risk picture into an auditable operational report — and pushes it to the frontline in one click."
]])
set_runs(shape_by_name(a, "Text 5"), [
    ["Two channels:", " in-app inbox + WhatsApp Cloud API — tracked per channel"],
    ["Closed loop:", " read & acknowledge stream back to the district live"],
    ["Auditable:", " every alert is a Firestore document with delivery history"],
])
set_runs(shape_by_name(a, "Text 7"), [[
    "WhatsApp failure never blocks in-app delivery — the frontline always gets the message."
]])
a.shapes.add_picture("docs/screenshots/notify-center.png", Inches(0.72), Inches(1.92), height=Inches(4.55))
a.shapes.add_picture("docs/screenshots/field-inbox-card.png", Inches(3.78), Inches(1.92), width=Inches(2.72))

# Slide B: District PDF report
b = clone_template()
set_runs(shape_by_name(b, "Text 2"), [["One-Click District Report"]])
set_runs(shape_by_name(b, "Text 3"), [["The whole district, on paper"]])
set_runs(shape_by_name(b, "Text 4"), [[
    "One button in the command centre exports a full PDF situation report — district summary, all 15 facilities triaged most-urgent first, and per-facility stock-out forecasts."
]])
set_runs(shape_by_name(b, "Text 5"), [
    ["Same engines:", " it can never disagree with the map"],
    ["Meeting-ready:", " printable, shareable, boardroom-friendly"],
    ["Always current:", " generated on demand from live state"],
])
set_runs(shape_by_name(b, "Text 7"), [[
    "From live map to signed-off paperwork — the last mile of the decision loop."
]])
b.shapes.add_picture("docs/screenshots/report-page1.png", Inches(1.05), Inches(1.90), height=Inches(4.62))

# Slide C: Roadmap
c = clone_template()
set_runs(shape_by_name(c, "Text 1"), [["ROADMAP"]])
set_runs(shape_by_name(c, "Text 2"), [["Built for What Comes Next"]])
set_runs(shape_by_name(c, "Text 3"), [["Same engines, new reach"]])
set_runs(shape_by_name(c, "Text 4"), [[
    "Every item extends the deterministic engine pipeline that already runs — new signals and new surfaces, no rearchitecture."
]])
t5c = shape_by_name(c, "Text 5")
set_bullets_labeled(t5c, [
    ("Offline-first:", " store-and-forward updates + SMS channel"),
    ("Field View 2.0:", " worker-tested redesign of the frontline app"),
    ("Regional languages:", " each state rollout ships its own languages"),
    ("Readiness scores:", " new ops & facility metrics from new data"),
])
set_runs(shape_by_name(c, "Text 7"), [[
    "Sequenced by pilot feedback: connectivity first, then languages, then richer scores."
]])
c.shapes.add_picture("docs/screenshots/field-inbox.png", Inches(2.1), Inches(1.92), height=Inches(4.55))

# ---- Position: A/B after stress mode, C before the go-to-market slide ----
sldIdLst = pres.slides._sldIdLst
ids = list(sldIdLst)
ea, eb, ec = ids[-3], ids[-2], ids[-1]
for e in (ea, eb, ec):
    sldIdLst.remove(e)
sldIdLst.insert(TEMPLATE_IDX + 1, eb)
sldIdLst.insert(TEMPLATE_IDX + 1, ea)
sldIdLst.insert(11, ec)

pres.save(OUT)
print("WROTE", OUT, "slides:", len(pres.slides.__iter__.__self__._sldIdLst))
